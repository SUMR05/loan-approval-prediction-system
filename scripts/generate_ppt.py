"""
Builds a simple, presentation-ready PowerPoint deck for the Loan Approval
Prediction System: project overview, pipeline, dataset, models compared,
model comparison chart/table, confusion matrices, decision logic, Power BI
dashboard fields, and business insights.

Pulls real numbers from outputs/model_comparison.csv and
models/model_metadata.json, so the deck always reflects the last training
run - run train_models.py first if you want fresh numbers.

Run directly:
    python scripts/generate_ppt.py
"""

import json
import os

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
MODEL_COMPARISON_CSV = os.path.join(BASE_DIR, "outputs", "model_comparison.csv")
MODEL_METADATA_PATH = os.path.join(BASE_DIR, "models", "model_metadata.json")
CONFUSION_MATRIX_DIR = os.path.join(BASE_DIR, "outputs", "confusion_matrices")
OUTPUT_PPTX = os.path.join(BASE_DIR, "outputs", "Loan_Approval_Prediction_System.pptx")

# ---------------------------------------------------------------------------
# Design constants
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x1F, 0x38, 0x64)
TEAL = RGBColor(0x2E, 0x86, 0xAB)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
AMBER = RGBColor(0xF0, 0x9C, 0x1F)
RED = RGBColor(0xD9, 0x4C, 0x4C)
GREY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def add_header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.5))
        p2 = box.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.italic = True
        p2.font.color.rgb = GREY


def add_title_slide(prs, title, subtitle, footer):
    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.6), SLIDE_W, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.8))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.8))
    p2 = box2.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xCF, 0xDD, 0xEE)

    box3 = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5))
    p3 = box3.text_frame.paragraphs[0]
    p3.text = footer
    p3.font.size = Pt(13)
    p3.font.color.rgb = RGBColor(0x9F, 0xB3, 0xCC)
    return slide


def add_bullet_slide(prs, title, bullets, subtitle=None, bullet_font=20):
    slide = blank_slide(prs)
    add_header(slide, title, subtitle)
    top = Inches(1.9) if subtitle else Inches(1.5)
    box = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.7), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        indent = item.startswith("  ")
        p.text = ("•  " if not indent else "   -  ") + item.strip()
        p.font.size = Pt(bullet_font - 3 if indent else bullet_font)
        p.font.color.rgb = GREY if indent else RGBColor(0x22, 0x22, 0x22)
        p.space_after = Pt(10)
    return slide


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = blank_slide(prs)
    add_header(slide, title)

    def column(x, heading, items, color):
        box = slide.shapes.add_textbox(x, Inches(1.6), Inches(5.8), Inches(0.5))
        p = box.text_frame.paragraphs[0]
        p.text = heading
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = color

        body = slide.shapes.add_textbox(x, Inches(2.15), Inches(5.8), Inches(4.8))
        tf = body.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "•  " + item
            p.font.size = Pt(17)
            p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
            p.space_after = Pt(8)

    column(Inches(0.6), left_title, left_items, TEAL)
    column(Inches(6.9), right_title, right_items, AMBER)
    return slide


def add_table_slide(prs, title, headers, rows, highlight_row_idx=None):
    slide = blank_slide(prs)
    add_header(slide, title)

    n_rows, n_cols = len(rows) + 1, len(headers)
    left, top, width, height = Inches(0.8), Inches(1.9), Inches(11.7), Inches(0.7 * n_rows)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    for r, row in enumerate(rows, start=1):
        is_highlight = highlight_row_idx is not None and (r - 1) == highlight_row_idx
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF3, 0xEA) if is_highlight else LIGHT_BG
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(15)
            p.font.bold = is_highlight
            p.font.color.rgb = RGBColor(0x1B, 0x5E, 0x20) if is_highlight else RGBColor(0x22, 0x22, 0x22)
            p.alignment = PP_ALIGN.CENTER
    return slide


def add_chart_slide(prs, title, categories, series_dict):
    slide = blank_slide(prs)
    add_header(slide, title)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_dict.items():
        chart_data.add_series(name, values)

    x, y, cx, cy = Inches(1.0), Inches(1.7), Inches(11.3), Inches(5.4)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    colors = [TEAL, GREEN, AMBER, NAVY]
    for i, series in enumerate(chart.plots[0].series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colors[i % len(colors)]

    value_axis = chart.value_axis
    value_axis.minimum_scale = 0
    value_axis.maximum_scale = 1
    value_axis.has_major_gridlines = True

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(14)
    return slide


def add_images_slide(prs, title, image_paths, captions):
    slide = blank_slide(prs)
    add_header(slide, title)

    n = len(image_paths)
    margin = Inches(0.5)
    gap = Inches(0.3)
    total_width = SLIDE_W - 2 * margin - gap * (n - 1)
    img_width = Emu(int(total_width / n))
    top = Inches(1.8)

    for i, (path, caption) in enumerate(zip(image_paths, captions)):
        left = Emu(int(margin + i * (img_width + gap)))
        if os.path.exists(path):
            slide.shapes.add_picture(path, left, top, width=img_width)
        cap_box = slide.shapes.add_textbox(left, Inches(6.5), img_width, Inches(0.6))
        p = cap_box.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(15)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = NAVY
    return slide


def add_flow_slide(prs, title, steps, colors):
    slide = blank_slide(prs)
    add_header(slide, title)

    n = len(steps)
    box_w, box_h = Inches(2.9), Inches(1.5)
    gap = Inches(0.55)
    total_w = box_w * n + gap * (n - 1)
    start_x = int((SLIDE_W - total_w) / 2)
    y = Inches(2.9)

    for i, (label, color) in enumerate(zip(steps, colors)):
        x = Emu(int(start_x + i * (box_w + gap)))
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = WHITE
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        if i < n - 1:
            arrow_x = Emu(int(x + box_w))
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_x, Emu(int(y + box_h / 2 - Inches(0.15))), gap, Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GREY
            arrow.line.fill.background()
    return slide


# ---------------------------------------------------------------------------
# Build the deck
# ---------------------------------------------------------------------------
def main():
    comparison_df = pd.read_csv(MODEL_COMPARISON_CSV)
    with open(MODEL_METADATA_PATH, "r", encoding="utf-8") as f:
        metadata = json.load(f)
    best_model = metadata["best_model"]
    best_metrics = metadata["metrics"]

    prs = new_presentation()

    # 1. Title slide
    add_title_slide(
        prs,
        "Loan Approval Prediction System",
        "Predicting loan approval, amount, and rejection reasons using Machine Learning",
        "Python  |  SQL  |  Pandas  |  Scikit-learn  |  Power BI",
    )

    # 2. What is this project
    add_bullet_slide(
        prs,
        "What Is This Project?",
        [
            "A system that looks at a loan applicant's details and makes a real lending decision.",
            "It answers 3 questions:",
            "  1. Should the loan be Approved or Rejected?",
            "  2. If approved, how much amount should be given?",
            "  3. If rejected, what is the exact reason?",
            "Built end-to-end: raw data -> cleaning -> machine learning -> business decision -> dashboard.",
        ],
        bullet_font=22,
    )

    # 3. Tech stack
    add_two_column_slide(
        prs,
        "Technology Used",
        "Data & Modeling",
        ["Python - cleaning data & building the model", "Pandas - organizing & transforming data",
         "Scikit-learn - training & comparing ML models", "SQL (SQLite) - storing & querying applicant records"],
        "Reporting",
        ["Power BI - dashboards for approval rate, risk, income trends", "Joblib - saving the trained model for reuse",
         "Matplotlib/Seaborn - confusion matrix charts"],
    )

    # 4. Pipeline flow
    add_flow_slide(
        prs,
        "How It Works - The Pipeline",
        ["1. Clean\nData", "2. Load to\nSQL", "3. Engineer\nFeatures", "4. Train &\nCompare Models", "5. Predict &\nDecide"],
        [TEAL, NAVY, TEAL, NAVY, GREEN],
    )

    # 5. Dataset overview
    add_bullet_slide(
        prs,
        "Dataset Overview",
        [
            "5,000 sample loan applicant records",
            "Includes: Gender, Marital Status, Education, Employment Status, Income, Co-applicant Income",
            "Also includes: Loan Amount, Loan Term, Credit Score, Credit History, Existing Loans",
            "Target column: Loan_Status - Approved (Y) or Rejected (N)",
            "Data is realistic - includes missing values, duplicate entries, and inconsistent text, just like real bank data.",
        ],
    )

    # 6. Cleaning & feature engineering
    add_two_column_slide(
        prs,
        "Data Cleaning & Feature Engineering",
        "Data Cleaning",
        ["Fill missing values (mode for categories, median for numbers)",
         "Fix inconsistent text (spacing, capitalization)",
         "Remove duplicate applications",
         "Fix out-of-range values"],
        "New Features Created",
        ["Total Income (applicant + co-applicant)",
         "EMI - estimated monthly installment",
         "Balance Income - income left after EMI",
         "Debt-to-Income Ratio",
         "Risk Category - Low / Medium / High"],
    )

    # 7. Models used
    add_bullet_slide(
        prs,
        "Machine Learning Models Compared",
        [
            "Logistic Regression - simple, adds up weighted factors to decide. Easy to explain.",
            "Decision Tree - works like a flowchart of yes/no questions.",
            "Random Forest - builds many decision trees and lets them vote together.",
            "All 3 were trained on the same data and compared fairly using the same test set.",
        ],
        bullet_font=22,
    )

    # 8. Model comparison table
    headers = ["Model", "Accuracy", "Precision", "Recall", "F1-Score"]
    rows = []
    best_idx = None
    for i, r in comparison_df.iterrows():
        rows.append([r["model"], f"{r['accuracy']:.2f}", f"{r['precision']:.2f}", f"{r['recall']:.2f}", f"{r['f1_score']:.2f}"])
        if r["model"] == best_model:
            best_idx = i
    add_table_slide(prs, "Model Comparison - Results", headers, rows, highlight_row_idx=best_idx)

    # 9. Model comparison chart
    categories = list(comparison_df["model"])
    series_dict = {
        "Accuracy": list(comparison_df["accuracy"]),
        "Precision": list(comparison_df["precision"]),
        "Recall": list(comparison_df["recall"]),
        "F1-Score": list(comparison_df["f1_score"]),
    }
    add_chart_slide(prs, "Model Comparison - Chart", categories, series_dict)

    # 10. Confusion matrices
    image_names = ["confusion_matrix_logistic_regression.png", "confusion_matrix_decision_tree.png", "confusion_matrix_random_forest.png"]
    captions = ["Logistic Regression", "Decision Tree", "Random Forest"]
    image_paths = [os.path.join(CONFUSION_MATRIX_DIR, name) for name in image_names]
    add_images_slide(prs, "Confusion Matrices - Where Each Model Went Right or Wrong", image_paths, captions)

    # 11. Best model
    add_bullet_slide(
        prs,
        "Best Model Selected",
        [
            f"Winner: {best_model}",
            f"Accuracy: {best_metrics['accuracy']:.1%}   |   Precision: {best_metrics['precision']:.1%}   |   "
            f"Recall: {best_metrics['recall']:.1%}   |   F1-Score: {best_metrics['f1_score']:.1%}",
            "Selected using F1-Score, because it balances two types of mistakes:",
            "  Approving a risky applicant (bad for the bank)",
            "  Rejecting a good applicant (bad for the customer)",
            "Saved using joblib as best_model.pkl - ready to reuse without retraining.",
        ],
        bullet_font=22,
    )

    # 12. Decision logic flow
    add_flow_slide(
        prs,
        "From Prediction to Real Decision",
        ["ML Model\n(Approve/Reject)", "Affordability\nCheck", "Final Decision\n+ Amount + Reason"],
        [TEAL, AMBER, GREEN],
    )
    add_bullet_slide(
        prs,
        "Loan Decision Logic (Simple Explanation)",
        [
            "Step 1: The ML model predicts if the applicant is creditworthy (Approve / Reject).",
            "Step 2: A simple rule checks how much loan the applicant can actually afford, based on income.",
            "Result: Approved (full amount), Partially Approved (reduced amount), or Rejected.",
            "If rejected, the system lists plain-language reasons - low credit score, high debt, unemployment, etc.",
            "All amounts are shown in Rupees (e.g. Rs 1,20,000).",
        ],
    )

    # 13. Power BI dashboard
    add_bullet_slide(
        prs,
        "Power BI Dashboard",
        [
            "A single ready-to-use file: loan_dashboard_data.csv",
            "Dashboard shows: Approval rate, Rejection rate, Risk category breakdown",
            "Also shows: Credit score analysis, Income analysis, Loan amount trends",
            "Also shows: Applicant profile summary, Approved amount vs. requested amount",
            "Business teams can explore this without touching any code.",
        ],
    )

    # 14. Business insights
    add_bullet_slide(
        prs,
        "Business Insights",
        [
            "Credit history is the strongest predictor of loan approval.",
            "High debt-to-income ratio leads to rejection, even for small loan amounts.",
            "Unemployed applicants are rejected more often, even at similar income levels.",
            "The affordability check catches cases where a good applicant asks for too much money.",
            "Outcome: faster, consistent, explainable loan decisions - less manual work for loan officers.",
        ],
    )

    # 15. Conclusion
    add_bullet_slide(
        prs,
        "Conclusion",
        [
            "Built a complete, working loan approval system - from raw data to final decision.",
            "Compared 3 machine learning models and picked the best one using proper evaluation.",
            "Added a real-world decision layer: approved amount + rejection reason, not just yes/no.",
            "Delivered a live Power BI dashboard for business reporting.",
            "Fully deployable: saved model + preprocessing pipeline, ready to reuse.",
        ],
        bullet_font=22,
    )

    # 16. Thank you
    slide = add_title_slide(prs, "Thank You", "Questions?", "Loan Approval Prediction System")

    prs.save(OUTPUT_PPTX)
    print(f"Presentation saved -> {OUTPUT_PPTX}")
    print(f"Slides: {len(prs.slides._sldIdLst)}")


if __name__ == "__main__":
    main()
